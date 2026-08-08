#!/usr/bin/env python3
"""Generate the NIRACONCHEM AI architecture deck.

Run:  python build_ppt.py
Output: NIRACONCHEM_AI_Architecture.pptx  (in project root)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY      = RGBColor(0x0B, 0x12, 0x20)   # background
NAVY2     = RGBColor(0x12, 0x1E, 0x33)   # card
NAVY3     = RGBColor(0x1B, 0x2A, 0x45)   # lighter card
SLATE     = RGBColor(0x9F, 0xB2, 0xCC)   # muted text
LIGHT     = RGBColor(0xEC, 0xF2, 0xFA)   # body text
TEAL      = RGBColor(0x22, 0xD3, 0xEE)   # primary accent
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # secondary accent
GREEN     = RGBColor(0x34, 0xD3, 0x99)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
RED       = RGBColor(0xF8, 0x70, 0x70)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    # de-select
    sp = bg._element.spPr
    sp.find(qn('a:lstStyle'))  # noop
    return s

def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, round_=True):
    shp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, sz, b, col, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.italic = it
            r.font.color.rgb = col; r.font.name = FONT
    return tb

def title_bar(s, kicker, title, accent=TEAL):
    rect(s, 0, 0, SW, Inches(1.15), NAVY2)
    rect(s, 0, 0, Inches(0.18), Inches(1.15), accent)
    txt(s, Inches(0.55), Inches(0.12), Inches(12), Inches(0.35),
        [[(kicker.upper(), 12, True, accent)]])
    txt(s, Inches(0.55), Inches(0.42), Inches(12.2), Inches(0.66),
        [[(title, 26, True, LIGHT)]], anchor=MSO_ANCHOR.MIDDLE)

def footer(s, n):
    txt(s, Inches(0.55), Inches(7.05), Inches(9), Inches(0.3),
        [[("NIRACONCHEM AI  —  Architecture Overview", 9, False, SLATE)]])
    txt(s, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
        [[(f"{n:02d}", 9, True, SLATE)]], align=PP_ALIGN.RIGHT)

def bullet(s, x, y, w, h, items, size=13, gap=7, accent=TEAL, head=None):
    runs = []
    if head:
        runs.append([(head, size+1, True, LIGHT)])
    for it in items:
        if isinstance(it, tuple):
            label, body = it
            runs.append([("▸ ", size, True, accent), (label, size, True, LIGHT),
                         (body, size, False, SLATE)])
        else:
            runs.append([("▸ ", size, True, accent), (it, size, False, SLATE)])
    txt(s, x, y, w, h, runs, space_after=gap)

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.5), SW, Inches(0.05), TEAL)
rect(s, 0, Inches(4.95), SW, Inches(0.05), AMBER)
txt(s, Inches(1.0), Inches(2.75), Inches(11.3), Inches(1.4),
    [[("NIRACONCHEM AI", 54, True, LIGHT)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.02), Inches(4.0), Inches(11.3), Inches(0.9),
    [[("Construction-Chemical Recommendation Platform", 24, False, TEAL)],
     [("How the AI works — a clear architecture walkthrough", 16, False, SLATE)]],
    space_after=6)
txt(s, Inches(1.02), Inches(5.2), Inches(11.3), Inches(1.2),
    [[("A product-profile-first assistant for UAE / GCC site teams.", 14, False, LIGHT)],
     [("Founded by Sravani Uppu — 10 years of construction-chemicals specification experience.", 13, False, SLATE)]],
    space_after=4)
txt(s, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
    [[("Next.js · FastAPI · LangGraph · Groq · Local JSON RAG", 12, True, AMBER)]])

# =====================================================================
# SLIDE 2 — PROBLEM & PLATFORM
# =====================================================================
s = slide(); title_bar(s, "Why this exists", "The Problem & The Platform")
# left card - problem
rect(s, Inches(0.55), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2,
     line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.7), Inches(5.4), Inches(0.5),
    [[("THE CHALLENGE", 14, True, RED)]])
bullet(s, Inches(0.8), Inches(2.25), Inches(5.45), Inches(3.7), [
    "UAE / GCC sites face extreme heat, UV, chloride, coastal & hydrostatic exposure.",
    "Choosing the wrong chemical system causes leaks, failures & rework.",
    "Generic AI chatbots answer from general knowledge — not grounded in real product data.",
    "Site teams need a correct system fast, with datasheet-backed reasoning.",
], size=13, gap=10)
# right card - platform
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2,
     line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.7), Inches(5.4), Inches(0.5),
    [[("THE PLATFORM", 14, True, GREEN)]])
bullet(s, Inches(7.1), Inches(2.25), Inches(5.45), Inches(3.7), [
    "Product-profile-first: recommendations pulled from a structured product dataset.",
    "LangGraph routes greetings, brand, general, knowledge & technical queries differently.",
    "Direct answers first — no long clarification forms before helping.",
    "Datasheet-backed RAG, file analysis, and a generated PDF technical report.",
    "UAE/GCC-aware: heat, UV, chloride, coastal, traffic & wet-area logic.",
], size=13, gap=10)
footer(s, 2)

# =====================================================================
# SLIDE 3 — END-TO-END ARCHITECTURE (flow)
# =====================================================================
s = slide(); title_bar(s, "Big picture", "End-to-End Architecture")
nodes = [
    ("User / Site Team", "Types a requirement\nor uploads a file", TEAL),
    ("Next.js Frontend", "Chat UI · File upload\nMarket · PDF · PWA", TEAL),
    ("FastAPI Backend", "/chat · /recommend\n/analyze-file · /rag", PURPLE),
    ("LangGraph Agent", "Intent routing\n& orchestration", AMBER),
    ("Knowledge + LLM", "Product profiles\nRAG chunks · Groq", GREEN),
]
n = len(nodes)
x0 = Inches(0.55); total_w = Inches(12.2)
gap = Inches(0.35)
bw = Emu(int((total_w - gap*(n-1)) / n))
y = Inches(2.3); bh = Inches(1.5)
for i,(t,sub,col) in enumerate(nodes):
    x = Emu(int(x0) + i*(int(bw)+int(gap)))
    rect(s, x, y, bw, bh, NAVY2, line=col)
    rect(s, x, y, bw, Inches(0.12), col)
    txt(s, x, Emu(int(y)+int(Inches(0.22))), bw, Inches(0.5),
        [[(t, 13, True, LIGHT)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(y)+int(Inches(0.72))), bw, Inches(0.7),
        [[(line, 9.5, False, SLATE)] for line in sub.split("\n")],
        align=PP_ALIGN.CENTER, space_after=0)
    if i < n-1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                 Emu(int(x)+int(bw)+Emu(int(gap))//2 - int(Inches(0.12))),
                 Emu(int(y)+int(bh)//2-int(Inches(0.12))),
                 Inches(0.24), Inches(0.24))
        ar.fill.solid(); ar.fill.fore_color.rgb = SLATE; ar.line.fill.background()
        ar.shadow.inherit=False
# return path
rect(s, x0, Inches(4.35), total_w, Inches(0.04), RGBColor(0x33,0x4d,0x73))
txt(s, x0, Inches(4.5), total_w, Inches(0.35),
    [[("RESPONSE PATH  →  Direct chat answer  ·  Market result cards  ·  PDF technical report", 12, True, AMBER)]],
    align=PP_ALIGN.CENTER)
# data layer
rect(s, Inches(0.55), Inches(5.1), Inches(12.2), Inches(1.5), NAVY3,
     line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.4),
    [[("DATA LAYER  (local JSON — fast, simple, swappable to Postgres + vector DB)", 13, True, TEAL)]])
txt(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.8),
    [[("product_profiles.json", 12, True, LIGHT), ("  — 24 structured product/system records   ", 12, False, SLATE),
      ("rag_chunks.json", 12, True, LIGHT), ("  — datasheet chunks   ", 12, False, SLATE),
      ("rag_index.json", 12, True, LIGHT), ("  — indexing metadata", 12, False, SLATE)]])
footer(s, 3)

# =====================================================================
# SLIDE 4 — FRONTEND
# =====================================================================
s = slide(); title_bar(s, "Layer 1", "Frontend — Next.js · React · TypeScript")
rect(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.4), [[("RESPONSIBILITIES", 13, True, TEAL)]])
bullet(s, Inches(0.8), Inches(2.2), Inches(5.6), Inches(3.8), [
    "Chat interface for construction-chemical questions",
    ("File upload", " PDF / DOCX / XLSX / TXT project analysis"),
    ("Market result", " product cards ranked for the query"),
    ("PDF download", " when report inputs are complete"),
    "Dark / light theme toggle (liquid-glass styling)",
    "Installable PWA (manifest + service worker)",
    "Mobile-responsive chat layout",
], size=12.5, gap=7)
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.7), Inches(5.4), Inches(0.4), [[("SMART CLIENT LOGIC", 13, True, AMBER)]])
bullet(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(3.8), [
    ("isConstructionRelatedQuery", " gates project-mode vs general"),
    ("tokenizeMarketText + scoreMarketProduct", " rank QCON products"),
    ("normalizeChatReply", " guard that keeps genuine RAG answers verbatim, only rewrites the legacy canned template"),
    ("apiFetch", " 45s timeout, talks to /chat, /analyze-file, /recommend/report"),
    ("Key files", " page.tsx, globals.css, data/qcon-market-products.json, public/ assets"),
], size=12, gap=6)
footer(s, 4)

# =====================================================================
# SLIDE 5 — BACKEND
# =====================================================================
s = slide(); title_bar(s, "Layer 2", "Backend — FastAPI service")
rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.5), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.4), [[("app/main.py — single service, clear modules", 14, True, TEAL)]])
bullet(s, Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.8), [
    "CORS restricted to configured frontend origins + *.vercel.app",
    "ensure_rag_indexes() auto-builds indexes on first request if missing",
    "Rule-based recommendation engine + optional Groq LLM enhancement",
    "PDF generation via ReportLab; session memory in chat_sessions",
], size=12, gap=4)
cards = [
    ("chat_agent_langgraph.py", "LangGraph chat brain + graph", PURPLE),
    ("rag_store.py", "Profile load, RAG retrieve, scoring", TEAL),
    ("rag_ingest.py", "Datasheet → chunks + profiles", GREEN),
    ("file_parser.py", "PDF/DOCX/XLSX/TXT extraction", AMBER),
    ("chat_sessions.py", "In-memory session store", RED),
    ("agent_prompt.py", "System prompt for the LLM", SLATE),
]
cw = Inches(3.9); ch = Inches(1.35); gx = Inches(0.3); gy = Inches(0.3)
for i,(t,sub,col) in enumerate(cards):
    r = i//3; c = i%3
    x = Emu(int(Inches(0.55)) + c*(int(cw)+int(gx)))
    y = Emu(int(Inches(3.3)) + r*(int(ch)+int(gy)))
    rect(s, x, y, cw, ch, NAVY2, line=col)
    rect(s, x, y, Inches(0.1), ch, col)
    txt(s, Emu(int(x)+int(Inches(0.25))), Emu(int(y)+int(Inches(0.18))), Emu(int(cw)-int(Inches(0.3))), Inches(0.5),
        [[(t, 12.5, True, LIGHT)]])
    txt(s, Emu(int(x)+int(Inches(0.25))), Emu(int(y)+int(Inches(0.62))), Emu(int(cw)-int(Inches(0.3))), Inches(0.6),
        [[(sub, 11, False, SLATE)]])
footer(s, 5)

# =====================================================================
# SLIDE 6 — LANGGRAPH CHAT BRAIN
# =====================================================================
s = slide(); title_bar(s, "The core", "LangGraph Chat Brain — the agent graph")
states = [
    ("START", "entry", SLATE),
    ("node_normalise", "clean + build\ncontext/history", TEAL),
    ("node_route_intent", "pick intent", AMBER),
]
# top chain
y = Inches(1.6)
prev = None
chain = [
    ("START", SLATE, "entry"),
    ("node_normalise", TEAL, "clean + join history"),
    ("node_route_intent", AMBER, "classify intent"),
]
bx = Inches(0.55); bw2 = Inches(2.55); bh2 = Inches(1.0); gapx = Inches(0.55)
for i,(t,col,sub) in enumerate(chain):
    x = Emu(int(bx) + i*(int(bw2)+int(gapx)))
    rect(s, x, y, bw2, bh2, NAVY2, line=col)
    txt(s, x, Emu(int(y)+int(Inches(0.12))), bw2, Inches(0.4), [[(t, 12, True, LIGHT)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(y)+int(Inches(0.5))), bw2, Inches(0.45), [[(sub, 9.5, False, SLATE)]], align=PP_ALIGN.CENTER)
    if i>0:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
              Emu(int(bx)+int(bw2)+(i-1)*(int(bw2)+int(gapx)) + int(gapx)//2 - int(Inches(0.11))),
              Emu(int(y)+int(bh2)//2-int(Inches(0.11))), Inches(0.22), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb=SLATE; ar.line.fill.background(); ar.shadow.inherit=False
# branch arrow down
centerx = Emu(int(bx) + 2*(int(bw2)+int(gapx)) + int(bw2)//2)
ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(int(centerx)-int(Inches(0.11))),
        Emu(int(y)+int(bh2)+int(Inches(0.05))), Inches(0.22), Inches(0.3))
ar.fill.solid(); ar.fill.fore_color.rgb=AMBER; ar.line.fill.background(); ar.shadow.inherit=False
# intent fan-out
intents = [
    ("brand_identity", "Founder/owner", RED),
    ("greeting", "Hi / hello", GREEN),
    ("general_question", "Casual chat", SLATE),
    ("knowledge_question", "“what is…?”", TEAL),
    ("technical_consultation", "Product / system", AMBER),
]
iy = Inches(3.35); iw = Inches(2.35); ih = Inches(0.95); igx = Inches(0.27)
for i,(t,sub,col) in enumerate(intents):
    x = Emu(int(Inches(0.4)) + i*(int(iw)+int(igx)))
    rect(s, x, iy, iw, ih, NAVY2, line=col)
    txt(s, x, Emu(int(iy)+int(Inches(0.12))), iw, Inches(0.4), [[(t, 11.5, True, LIGHT)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(iy)+int(Inches(0.5))), iw, Inches(0.4), [[(sub, 9.5, False, SLATE)]], align=PP_ALIGN.CENTER)
# merge to recommend
ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(int(centerx)-int(Inches(0.11))),
        Emu(int(iy)+int(ih)+int(Inches(0.05))), Inches(0.22), Inches(0.3))
ar.fill.solid(); ar.fill.fore_color.rgb=AMBER; ar.line.fill.background(); ar.shadow.inherit=False
ry = Inches(4.95)
rect(s, Emu(int(centerx)-int(Inches(1.4))), ry, Inches(2.8), Inches(0.95), NAVY2, line=AMBER)
txt(s, Emu(int(centerx)-int(Inches(1.4))), Emu(int(ry)+int(Inches(0.12))), Inches(2.8), Inches(0.4),
    [[("node_recommend", 12.5, True, LIGHT)]], align=PP_ALIGN.CENTER)
txt(s, Emu(int(centerx)-int(Inches(1.4))), Emu(int(ry)+int(Inches(0.5))), Inches(2.8), Inches(0.4),
    [[("build answer → END", 10, False, SLATE)]], align=PP_ALIGN.CENTER)
ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(int(centerx)-int(Inches(0.11))),
        Emu(int(ry)+int(Inches(0.95))+int(Inches(0.05))), Inches(0.22), Inches(0.3))
ar.fill.solid(); ar.fill.fore_color.rgb=GREEN; ar.line.fill.background(); ar.shadow.inherit=False
txt(s, Emu(int(centerx)-int(Inches(1.0))), Inches(6.45), Inches(2.0), Inches(0.4),
    [[("END", 12, True, GREEN)]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.55), Inches(6.85), Inches(12), Inches(0.4),
    [[("Technical & knowledge paths first ", 11, False, SLATE),
      ("retrieve product profiles + datasheet chunks", 11, True, TEAL),
      (" then answer via Groq (or grounded rule fallback).", 11, False, SLATE)]], align=PP_ALIGN.CENTER)
footer(s, 6)

# =====================================================================
# SLIDE 7 — INTENT ROUTING LOGIC
# =====================================================================
s = slide(); title_bar(s, "Routing", "How the Intent Is Decided")
rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.4),
    [[("node_route_intent — ordered checks (first match wins):", 14, True, TEAL)]])
checks = [
    ("1. Brand terms", "founder / founded / owner / “niraconchem” / “sravani” → brand_identity"),
    ("2. Greeting", "hi/hello + ≤5 words + no history → greeting"),
    ("3. Product-profile match", "query matches a known product/brand name → technical_consultation"),
    ("4. Knowledge shape", "ends with “?” or starts with “what is / explain / define…” → general_question"),
    ("5. Slot detail", "contains a requirement / area / substrate / exposure / location term → technical_consultation"),
    ("6. Construction score", "construction_score(query + ½ context) ≥ 2 → technical_consultation"),
    ("7. Default", "otherwise → general_question (LLM answers naturally)"),
]
y = Inches(2.25)
for (h,b) in checks:
    rect(s, Inches(0.8), y, Inches(0.18), Inches(0.55), AMBER)
    txt(s, Inches(1.15), y, Inches(11.3), Inches(0.6),
        [[(h+"  —  ", 13, True, LIGHT),(b, 12.5, False, SLATE)]], anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y)+int(Inches(0.6)))
footer(s, 7)

# =====================================================================
# SLIDE 8 — RECOMMENDATION ENGINE / SCORING
# =====================================================================
s = slide(); title_bar(s, "Brainpower", "Recommendation Engine — Retrieval & Scoring")
rect(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.68), Inches(5.6), Inches(0.4), [[("retrieve_product_profiles()", 13, True, TEAL)]])
bullet(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(3.8), [
    "Tokenize query, drop stopwords, expand synonyms (roof→roofing,membrane,uv)",
    "Score each profile = token overlap ÷ √profile_tokens",
    ("Boosts", " +0.45 area match, +0.35 category, +0.55 requested area, +0.25 known manufacturer"),
    ("Penalties", " −0.8 EIFS when not wanted, −0.3 roof when not asked"),
    "Dedupe by (product, brand); return top-N",
], size=12, gap=7)
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.68), Inches(5.4), Inches(0.4), [[("profile_score() — per-product weight", 13, True, AMBER)]])
bullet(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(3.8), [
    ("Exact product name in query", " +220"),
    ("Full name-token coverage", " +170 / +120 / +60 by ratio"),
    ("Brand + name both match", " +35"),
    ("Phrase boosts", " “tile adhesive”, “waterproofing”, “concrete repair” +35"),
    ("“why / how / where” + 2 name tokens", " +70"),
    ("Has price / url / description / brand", " small + bonuses"),
    "Result: ranked, grounded product/system matches",
], size=11.5, gap=5)
footer(s, 8)

# =====================================================================
# SLIDE 9 — DATA / RAG LAYER
# =====================================================================
s = slide(); title_bar(s, "Knowledge", "Data & RAG Layer")
left = [
    ("product_profiles.json", "24 structured product/system records", TEAL),
    ("rag_chunks.json", "datasheet chunks (900-word, 160 overlap)", GREEN),
    ("rag_index.json", "document profile metadata", AMBER),
]
x = Inches(0.55)
for (t,sub,col) in left:
    rect(s, x, Inches(1.6), Inches(3.95), Inches(1.25), NAVY2, line=col)
    txt(s, Emu(int(x)+int(Inches(0.2))), Inches(1.75), Inches(3.6), Inches(0.5), [[(t, 12, True, LIGHT)]])
    txt(s, Emu(int(x)+int(Inches(0.2))), Inches(2.2), Inches(3.6), Inches(0.55), [[(sub, 10.5, False, SLATE)]])
    x = Emu(int(x)+int(Inches(4.1)))
rect(s, Inches(0.55), Inches(3.1), Inches(12.2), Inches(3.0), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(3.22), Inches(11.7), Inches(0.4), [[("rag_ingest.py — datasheet → intelligence", 14, True, TEAL)]])
bullet(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(2.3), [
    ("Sources", " 13 manufacturer .docx datasheets in data/datasheets/"),
    ("Chunking", " 900 words, 160-word overlap, each stored with tokens + keywords"),
    ("Profile build", " detect manufacturer/country, system type, layers (primer/membrane/reinforcement/top-coat), application areas, performance (bond/tensile/elongation/SRI…)"),
    ("Product patterns", " regex for Saveto / Vetoprime / Vetoproof / Vetotop families"),
    ("System rules", " map text → PVC roofing, Acrylic-PU hybrid, epoxy coating, repair mortar, sealant…"),
    ("Retrieval", " rag_store scores chunks by overlap + area/roof/repair boosts, diverse by source file"),
], size=12, gap=5)
footer(s, 9)

# =====================================================================
# SLIDE 10 — UAE/GCC CLIMATE & EXPOSURE RULES
# =====================================================================
s = slide(); title_bar(s, "Domain logic", "UAE / GCC Climate & Exposure Rules")
rect(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.68), Inches(5.6), Inches(0.4), [[("CLIMATE_RULES", 13, True, AMBER)]])
bullet(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(3.8), [
    ("Coastal", " humidity + chloride → UV-stable, salt-tolerant systems"),
    ("Desert", " heat/dry wind → hot-weather windows, moisture loss control"),
    ("Rooftop", " UV + thermal movement → flexible UV-resistant membranes"),
    ("Underground", " hydrostatic pressure → strong substrate prep systems"),
    ("Industrial", " chemical/abrasion → confirm resistance before selection"),
], size=12, gap=8)
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.68), Inches(5.4), Inches(0.4), [[("SLOT / CATEGORY LOGIC", 13, True, TEAL)]])
bullet(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(3.8), [
    ("UAE emirates", " mapped; coastal vs desert exposure hints"),
    ("CATEGORY_RULES", " waterproof/repair/tile/floor/corrosion/sealant → candidate systems"),
    ("5 intake slots", " problem, area, substrate, exposure, location"),
    ("infer_missing_information", " flags what’s still needed for a report"),
    ("Caveat", " output is technical guidance — verify vs manufacturer datasheet & spec"),
], size=12, gap=8)
footer(s, 10)

# =====================================================================
# SLIDE 11 — FILE UPLOAD & ANALYSIS
# =====================================================================
s = slide(); title_bar(s, "Inputs", "File Upload & Project Analysis")
flow = [
    ("Upload", "PDF / DOCX\nXLSX / TXT", TEAL),
    ("file_parser", "extract text\n(12k char cap)", PURPLE),
    ("summarize signals", "locations\nareas · requirements", AMBER),
    ("/analyze-file", "returns JSON\n+ preview", GREEN),
]
x = Inches(0.55); bw3=Inches(2.85); bh3=Inches(1.3); gp=Inches(0.45)
for i,(t,sub,col) in enumerate(flow):
    xc = Emu(int(x) + i*(int(bw3)+int(gp)))
    rect(s, xc, Inches(1.8), bw3, bh3, NAVY2, line=col)
    txt(s, xc, Emu(int(Inches(1.8))+int(Inches(0.15))), bw3, Inches(0.45), [[(t,13,True,LIGHT)]], align=PP_ALIGN.CENTER)
    txt(s, xc, Emu(int(Inches(1.8))+int(Inches(0.6))), bw3, Inches(0.6),
        [[(l,9.5,False,SLATE)] for l in sub.split("\n")], align=PP_ALIGN.CENTER, space_after=0)
    if i<len(flow)-1:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
             Emu(int(xc)+int(bw3)+int(gp)//2-int(Inches(0.11))),
             Inches(2.4), Inches(0.22), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb=SLATE; ar.line.fill.background(); ar.shadow.inherit=False
rect(s, Inches(0.55), Inches(3.6), Inches(12.2), Inches(2.5), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(3.72), Inches(11.7), Inches(0.4), [[("What it extracts", 14, True, TEAL)]])
bullet(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.8), [
    ("LOCATION_KEYWORDS", " Dubai, Abu Dhabi, Sharjah, RAK, Fujairah, Al Ain…"),
    ("AREA_KEYWORDS", " basement, roof, pool, bathroom, foundation, facade, parking, slab…"),
    ("REQUIREMENT_KEYWORDS", " waterproofing, repair, tile adhesive, grout, sealant, coating, crack…"),
    ("Feeds forward", " extracted context is sent into /chat & /recommend/report to ground the answer"),
], size=12.5, gap=7)
footer(s, 11)

# =====================================================================
# SLIDE 12 — PDF REPORT
# =====================================================================
s = slide(); title_bar(s, "Output", "PDF Technical Report Generation")
rect(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.68), Inches(5.6), Inches(0.4), [[("/recommend/report", 14, True, AMBER)]])
bullet(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(3.8), [
    "Built with ReportLab (tables, styles, images, page breaks)",
    ("report_ready", " true only when slots are complete: problem, area, substrate, exposure, location"),
    ("Groq enhancement", " strict JSON: best_system, manufacturer, products, why_recommended, precautions, questions"),
    ("Priority order", " product profile → RAG datasheet context → LLM fills gaps only"),
    "Frontend shows the PDF download action when ready",
], size=12, gap=8)
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.68), Inches(5.4), Inches(0.4), [[("Report contains", 14, True, TEAL)]])
bullet(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(3.8), [
    "Project summary & detected location",
    "Climate / exposure context",
    "Recommended category & application guidance",
    "Selected product profile + alternatives",
    "Why recommended + supporting datasheet references",
    "Precautions & any open questions",
], size=12.5, gap=8)
footer(s, 12)

# =====================================================================
# SLIDE 13 — MARKET RESULT
# =====================================================================
s = slide(); title_bar(s, "Discovery", "Market Result — QCON Products")
rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.0), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.4), [[("Frontend-side product discovery", 14, True, TEAL)]])
bullet(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.3), [
    ("Source", " qcon-market-products.json (scraped QCON product data) + local images in public/assets/qcon-products/"),
    ("Trigger", " when intent = technical_consultation and no clarification needed, UI switches to the Market tab"),
    ("Ranking", " scoreMarketProduct() tokenizes query + reply + captured requirements, weights name/category/brand/keyword/description"),
], size=12, gap=5)
rect(s, Inches(0.55), Inches(3.8), Inches(12.2), Inches(2.3), NAVY3, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.4), [[("Each product card", 14, True, AMBER)]])
txt(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.5),
    [[("name · company · brand · category · price · url · imageUrl · localImage · description · keywords", 12.5, True, LIGHT)],
     [("→ Shows practical, purchaseable options next to the AI recommendation so users can act immediately.", 12.5, False, SLATE)]],
    space_after=8)
footer(s, 13)

# =====================================================================
# SLIDE 14 — API ENDPOINTS
# =====================================================================
s = slide(); title_bar(s, "Interface", "API Endpoints (FastAPI, :8000)")
eps = [
    ("GET  /health", "backend health check", GREEN),
    ("POST /chat", "LangGraph chat response", TEAL),
    ("POST /recommend", "rule-based recommendation", AMBER),
    ("POST /recommend/report", "PDF report download", AMBER),
    ("POST /analyze-file", "upload + analyze project files", PURPLE),
    ("POST /rag/ingest", "rebuild RAG/product indexes", RED),
    ("GET  /rag/status", "inspect RAG index readiness", RED),
]
x = Inches(0.55); cw=Inches(3.95); ch=Inches(0.95); gx=Inches(0.2); gy=Inches(0.25)
for i,(t,sub,col) in enumerate(eps):
    r=i//3; c=i%3
    xc = Emu(int(Inches(0.55))+c*(int(cw)+int(gx)))
    yc = Emu(int(Inches(1.6))+r*(int(ch)+int(gy)))
    rect(s, xc, yc, cw, ch, NAVY2, line=col)
    rect(s, xc, yc, Inches(0.1), ch, col)
    txt(s, Emu(int(xc)+int(Inches(0.25))), Emu(int(yc)+int(Inches(0.12))), Emu(int(cw)-int(Inches(0.3))), Inches(0.45),
        [[(t, 12.5, True, LIGHT)]])
    txt(s, Emu(int(xc)+int(Inches(0.25))), Emu(int(yc)+int(Inches(0.52))), Emu(int(cw)-int(Inches(0.3))), Inches(0.4),
        [[(sub, 11, False, SLATE)]])
footer(s, 14)

# =====================================================================
# SLIDE 15 — TECH STACK & DEPLOYMENT
# =====================================================================
s = slide(); title_bar(s, "Delivery", "Tech Stack & Deployment")
rect(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(0.8), Inches(1.68), Inches(5.6), Inches(0.4), [[("TECH STACK", 14, True, TEAL)]])
bullet(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(3.8), [
    ("Frontend", " Next.js, React, TypeScript, CSS, lucide-react"),
    ("Backend", " Python, FastAPI, Pydantic, LangGraph, Groq"),
    ("Documents", " pypdf, python-docx, openpyxl"),
    ("Reports", " ReportLab"),
    ("Data / RAG", " local JSON indexes, datasheet chunks, product profiles"),
    ("Deploy", " Vercel (frontend) + Render (backend)"),
], size=12.5, gap=8)
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.6), NAVY2, line=RGBColor(0x23,0x3A,0x5C))
txt(s, Inches(7.1), Inches(1.68), Inches(5.4), Inches(0.4), [[("ENVIRONMENT VARIABLES", 14, True, AMBER)]])
bullet(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(3.8), [
    ("Backend .env", " GROQ_API_KEY, GROQ_MODEL, FRONTEND_ORIGINS"),
    ("Frontend", " NEXT_PUBLIC_API_BASE_URL"),
    ("Render", " PYTHON_VERSION 3.11, GROQ_API_KEY (sync:false)"),
    ("Security", " keys never committed; CORS scoped; upload size limited"),
    ("Note", " works without Groq — grounded rule-based fallback answers"),
], size=12.5, gap=8)
footer(s, 15)

# =====================================================================
# SLIDE 16 — ROADMAP
# =====================================================================
s = slide(); title_bar(s, "Future", "Scaling Roadmap")
cols = [
    ("MOVE TO PRODUCTION", TEAL, [
        "Sessions: Redis / Postgres",
        "Profiles: Postgres + admin UI",
        "RAG: Qdrant / Pinecone / Chroma",
    ]),
    ("RELIABILITY", AMBER, [
        "Auth + role-based access",
        "Observability: logs, traces, evals",
        "Background workers (ingest, PDF)",
    ]),
    ("INTELLIGENCE", PURPLE, [
        "Camera / site-image upload",
        "Vision: cracks, dampness, defects",
        "Manufacturer comparison engine",
    ]),
    ("REACH", GREEN, [
        "Multi-language (UAE/GCC)",
        "User accounts & workspaces",
        "Live climate / weather API",
    ]),
]
cw=Inches(2.95); gx=Inches(0.2)
for i,(h,col,items) in enumerate(cols):
    x = Emu(int(Inches(0.55))+i*(int(cw)+int(gx)))
    rect(s, x, Inches(1.6), cw, Inches(4.4), NAVY2, line=col)
    rect(s, x, Inches(1.6), cw, Inches(0.55), col)
    txt(s, x, Inches(1.68), cw, Inches(0.4), [[(h, 12.5, True, NAVY if col!=SLATE else LIGHT)]], align=PP_ALIGN.CENTER)
    yy = Inches(2.35)
    for it in items:
        rect(s, Emu(int(x)+int(Inches(0.2))), yy, Inches(0.12), Inches(0.5), col)
        txt(s, Emu(int(x)+int(Inches(0.45))), yy, Emu(int(cw)-int(Inches(0.6))), Inches(0.9),
            [[(it, 11.5, False, LIGHT)]], anchor=MSO_ANCHOR.MIDDLE)
        yy = Emu(int(yy)+int(Inches(1.15)))
footer(s, 16)

# =====================================================================
OUT = "NIRACONCHEM_AI_Architecture.pptx"
prs.save(OUT)
print("Saved", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
